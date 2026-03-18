import json
from pptx import Presentation

from services.llm_client import cerebras_chat


ROLE_SLIDE_GUIDANCE = {
    "teacher": """
Audience style:
- Create slides for teaching and explanation.
- Focus on clarity, learning flow, and easy understanding.
- Prefer concept, working, example, and takeaway structure.

Presentation style:
- Bullets should read like teaching points.
- Keep terms simple and self-explanatory.
- Maintain a step-by-step order across slides.
""",
    "student": """
Audience style:
- Create slides for a student presenting a final year project.
- Focus on problem, objective, method, modules, output, and conclusion.
- Make the content confident, practical, and viva-ready.

Presentation style:
- Bullets should support spoken explanation during project review.
- Include implementation and result details where useful.
- Keep the tone academic and straightforward.
""",
    "company": """
Audience style:
- Create slides for a company or stakeholder presentation.
- Focus on problem, solution, workflow, impact, scalability, and results.
- Highlight business value and practical outcomes.

Presentation style:
- Bullets should be concise, polished, and business-oriented.
- Prefer impact-focused wording over classroom explanation.
- Keep the sequence executive-friendly and solution-oriented.
""",
}

def build_slide_prompt(text, slide_count, role="teacher"):
    role_guidance = ROLE_SLIDE_GUIDANCE.get(role, ROLE_SLIDE_GUIDANCE["teacher"])
    return f"""
Create EXACTLY {slide_count} presentation slides from the content below.

Rules:
- Use all important information.
- Summarize clearly.
- Each slide must have:
  - One title
  - 4-6 bullet points
- The slide title and bullet points must always be in English.
- Even if the source content is Tamil or mixed-language, translate and present the slides in clear English.
- Match the slides to this presentation role:
{role_guidance}
- Output ONLY valid JSON in this format:

[
  {{
    "title": "Slide title",
    "bullets": ["point 1", "point 2"]
  }}
]

Content:
{text}
"""


def summarize_into_slides(text, slide_count, role="teacher"):
    chunks = chunk_text(text, max_chars=2500)
    partial_summaries = []

    for chunk in chunks:
        response = cerebras_chat.create_chat_completion(
            messages=[
                {
                    "role": "system",
                    "content": "Summarize the following content clearly and concisely.",
                },
                {
                    "role": "user",
                    "content": chunk,
                },
            ]
        )

        summary = response.choices[0].message.content.strip()
        partial_summaries.append(summary)

    combined_summary = "\n".join(partial_summaries)
    final_prompt = build_slide_prompt(combined_summary, slide_count, role=role)

    final_response = cerebras_chat.create_chat_completion(
        messages=[
            {"role": "system", "content": "You are an expert presentation designer."},
            {"role": "user", "content": final_prompt},
        ]
    )

    raw = final_response.choices[0].message.content.strip()
    return _extract_slides_json(raw)


def create_summary_ppt(slides_data, output_path):
    prs = Presentation()
    layout = prs.slide_layouts[1]

    for slide_data in slides_data:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data["title"]

        tf = slide.placeholders[1].text_frame
        tf.clear()

        for bullet in slide_data["bullets"]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 1

    prs.save(output_path)
    return output_path


def chunk_text(text, max_chars=3000):
    chunks = []
    current = ""

    for line in text.split("\n"):
        if len(current) + len(line) < max_chars:
            current += line + "\n"
        else:
            chunks.append(current)
            current = line + "\n"

    if current.strip():
        chunks.append(current)

    return chunks


def _extract_slides_json(raw_text):
    decoder = json.JSONDecoder()

    for start in (idx for idx, char in enumerate(raw_text) if char == "["):
        try:
            candidate, _ = decoder.raw_decode(raw_text[start:])
        except json.JSONDecodeError:
            continue

        if isinstance(candidate, list):
            return _normalize_slides(candidate)

    raise ValueError("Cerebras did not return a valid slide JSON array")


def _normalize_slides(slides):
    normalized = []

    for slide in slides:
        if not isinstance(slide, dict):
            continue

        title = str(slide.get("title", "")).strip()
        bullets = slide.get("bullets", [])

        if isinstance(bullets, str):
            bullets = [bullets]
        elif not isinstance(bullets, list):
            bullets = []

        clean_bullets = [str(item).strip() for item in bullets if str(item).strip()]
        if title and clean_bullets:
            normalized.append({"title": title, "bullets": clean_bullets})

    if not normalized:
        raise ValueError("Cerebras returned JSON, but it did not contain usable slides")

    return normalized
